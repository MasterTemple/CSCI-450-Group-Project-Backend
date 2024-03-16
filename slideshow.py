import re
from pptx import Presentation
from pptx.text.text import TextFrame, _Paragraph
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(hex_color: str) -> list[int]:
    return [int(c, 16) for c in re.findall(r"[A-z0-9]{2}", hex_color)]

BLANK_LAYOUT_ID = 6

class Slideshow:

    def __init__(self, slide_lines:list[list[str]], settings: dict, title: str, author: str) -> None:
        self.title = title
        self.author = author
        self.background_color = RGBColor(*hex_to_rgb(settings["backgroundColor"]))
        self.text_color = RGBColor(*hex_to_rgb(settings["textColor"]))
        self.font_size = Pt(settings["fontSize"])
        self.font_family = settings["fontFamily"]
        self.include_title = settings["includeTitleSlide"]
        self.presentation = Presentation()
        self.slide_lines = slide_lines
        self.settings = settings
        self.create_presentation()

    def create_paragraph(self, text: list | str, font_size: Pt):
        if type(text) is list:
            text = "\n".join(text)
        p: _Paragraph = self.tf.add_paragraph()
        p.line_spacing = 1.5
        p.space_before = Inches(0)
        p.space_after = Inches(0)
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = self.text_color
        p.font.name = self.font_family
        p.alignment = PP_ALIGN.CENTER

    def create_slide(self):
        # create blank slide
        self.slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[BLANK_LAYOUT_ID])
        # set background color
        self.slide.background.fill.solid()
        self.slide.background.fill.fore_color.rgb = self.background_color

    def create_text_box(self, height: int):
        # calculate center
        text_box_width = Inches(4)
        text_box_height = Inches(height)

        left = (self.presentation.slide_width / 2) - (text_box_width / 2)
        top = (self.presentation.slide_height / 2) - (text_box_height / 2)
        # create textbox in the center
        self.text_box = self.slide.shapes.add_textbox(left, top, text_box_width, text_box_height)

    def create_text_frame(self):
        self.tf: TextFrame = self.text_box.text_frame
    
    def remove_empty_first_line(self):
        self.tf.paragraphs[0]._element.getparent().remove(self.tf.paragraphs[0]._element)

    def create_presentation(self):
        if self.include_title:
            self.create_slide()
            self.create_text_box(2)
            self.create_text_frame()
            self.create_paragraph(self.title, Pt(54))
            self.create_paragraph(self.author, Pt(32))
            self.remove_empty_first_line()

        for lines in self.slide_lines:
            self.create_slide()
            self.create_text_box(len(lines))
            self.create_text_frame()
            self.create_paragraph(lines, self.font_size)
            self.remove_empty_first_line()

    # save the presentation
    def save(self, file: str):
        self.presentation.save(file)
